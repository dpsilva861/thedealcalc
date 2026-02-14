"""
Content extraction engine — reads the actual content of every file type.

Uses Apache Tika as the primary extractor with fallback to specialized
Python libraries for common formats. Supports:
- PDFs, Word, Excel, PowerPoint
- Images via OCR (Tesseract)
- Emails (.eml, .msg)
- Plain text and code files
- Audio/video metadata
- Archives (list contents)
"""

from __future__ import annotations

import email
import io
import json
import logging
import mimetypes
import os
import struct
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

# Extension to broad type mapping for dispatch
TEXT_EXTENSIONS = {
    ".txt", ".md", ".rst", ".csv", ".tsv", ".log", ".ini", ".cfg", ".conf",
    ".yaml", ".yml", ".json", ".xml", ".html", ".htm", ".css", ".js", ".ts",
    ".tsx", ".jsx", ".py", ".rb", ".php", ".java", ".c", ".cpp", ".h", ".hpp",
    ".cs", ".go", ".rs", ".swift", ".kt", ".sql", ".sh", ".bash", ".bat",
    ".ps1", ".psm1", ".toml", ".env", ".gitignore", ".dockerfile",
}

IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".tif", ".gif", ".webp"}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".flac", ".aac", ".ogg", ".wma", ".m4a"}
VIDEO_EXTENSIONS = {".mp4", ".avi", ".mkv", ".mov", ".wmv", ".flv", ".webm"}
ARCHIVE_EXTENSIONS = {".zip", ".rar", ".7z", ".tar", ".gz", ".bz2", ".xz"}


class ContentReader:
    """Extracts text content and metadata from files."""

    def __init__(self, tika_server_url: str = "", max_size: int = 100 * 1024 * 1024) -> None:
        self.tika_server_url = tika_server_url
        self.max_size = max_size
        self._tika_available: bool | None = None

    def extract(self, file_path: str) -> dict[str, Any]:
        """
        Extract content from a file. Returns dict with:
        - text: extracted text content
        - method: extraction method used
        - metadata: any additional metadata extracted
        - error: error message if extraction failed
        """
        result: dict[str, Any] = {
            "text": "",
            "method": "",
            "metadata": {},
            "error": "",
        }

        if not os.path.exists(file_path):
            result["error"] = "File not found"
            return result

        file_size = os.path.getsize(file_path)
        if file_size > self.max_size:
            result["error"] = f"File too large ({file_size} bytes)"
            return result

        if file_size == 0:
            result["error"] = "Empty file"
            return result

        ext = Path(file_path).suffix.lower()

        # Dispatch to the appropriate extractor
        try:
            if ext in TEXT_EXTENSIONS:
                result = self._extract_text_file(file_path)
            elif ext == ".pdf":
                result = self._extract_pdf(file_path)
            elif ext in (".doc", ".docx"):
                result = self._extract_word(file_path)
            elif ext in (".xls", ".xlsx"):
                result = self._extract_excel(file_path)
            elif ext in (".ppt", ".pptx"):
                result = self._extract_powerpoint(file_path)
            elif ext in (".eml",):
                result = self._extract_email(file_path)
            elif ext in (".msg",):
                result = self._extract_msg(file_path)
            elif ext in IMAGE_EXTENSIONS:
                result = self._extract_image(file_path)
            elif ext in AUDIO_EXTENSIONS or ext in VIDEO_EXTENSIONS:
                result = self._extract_media_metadata(file_path)
            elif ext in ARCHIVE_EXTENSIONS:
                result = self._extract_archive_listing(file_path)
            else:
                # Try Tika as a catch-all
                result = self._extract_with_tika(file_path)
                if not result["text"] and not result["error"]:
                    # Last resort: try reading as text
                    result = self._extract_text_file(file_path)
        except Exception as e:
            logger.warning("Extraction failed for %s: %s", file_path, e)
            result["error"] = str(e)

        return result

    def _extract_text_file(self, file_path: str) -> dict[str, Any]:
        """Read plain text files with encoding detection."""
        result: dict[str, Any] = {"text": "", "method": "text", "metadata": {}, "error": ""}
        try:
            # Try UTF-8 first
            with open(file_path, "r", encoding="utf-8") as f:
                result["text"] = f.read()
        except UnicodeDecodeError:
            try:
                import chardet
                with open(file_path, "rb") as f:
                    raw = f.read()
                detected = chardet.detect(raw)
                encoding = detected.get("encoding", "latin-1") or "latin-1"
                result["text"] = raw.decode(encoding, errors="replace")
                result["metadata"]["detected_encoding"] = encoding
            except Exception as e:
                result["error"] = f"Encoding detection failed: {e}"
        except Exception as e:
            result["error"] = str(e)
        return result

    def _extract_pdf(self, file_path: str) -> dict[str, Any]:
        """Extract text from PDF files."""
        result: dict[str, Any] = {"text": "", "method": "pypdf2", "metadata": {}, "error": ""}
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            pages = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
            result["text"] = "\n\n".join(pages)
            if reader.metadata:
                meta = reader.metadata
                result["metadata"] = {
                    "title": str(meta.title) if meta.title else "",
                    "author": str(meta.author) if meta.author else "",
                    "subject": str(meta.subject) if meta.subject else "",
                    "creator": str(meta.creator) if meta.creator else "",
                    "pages": len(reader.pages),
                }
        except ImportError:
            result = self._extract_with_tika(file_path)
        except Exception as e:
            result["error"] = str(e)
            # Fallback to Tika
            tika_result = self._extract_with_tika(file_path)
            if tika_result["text"]:
                return tika_result
        return result

    def _extract_word(self, file_path: str) -> dict[str, Any]:
        """Extract text from Word documents."""
        result: dict[str, Any] = {"text": "", "method": "python-docx", "metadata": {}, "error": ""}
        ext = Path(file_path).suffix.lower()
        if ext == ".docx":
            try:
                from docx import Document
                doc = Document(file_path)
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                # Also extract from tables
                for table in doc.tables:
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            paragraphs.append(" | ".join(cells))
                result["text"] = "\n".join(paragraphs)
                core = doc.core_properties
                result["metadata"] = {
                    "title": core.title or "",
                    "author": core.author or "",
                    "subject": core.subject or "",
                    "created": str(core.created) if core.created else "",
                    "modified": str(core.modified) if core.modified else "",
                }
            except ImportError:
                result = self._extract_with_tika(file_path)
            except Exception as e:
                result["error"] = str(e)
        else:
            # .doc format — use Tika
            result = self._extract_with_tika(file_path)
        return result

    def _extract_excel(self, file_path: str) -> dict[str, Any]:
        """Extract text from Excel files."""
        result: dict[str, Any] = {"text": "", "method": "openpyxl", "metadata": {}, "error": ""}
        ext = Path(file_path).suffix.lower()
        if ext == ".xlsx":
            try:
                from openpyxl import load_workbook
                wb = load_workbook(file_path, read_only=True, data_only=True)
                sheets_text = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    sheet_lines = [f"[Sheet: {sheet_name}]"]
                    for row in ws.iter_rows(values_only=True):
                        cells = [str(c) if c is not None else "" for c in row]
                        if any(c.strip() for c in cells):
                            sheet_lines.append(" | ".join(cells))
                    sheets_text.append("\n".join(sheet_lines))
                result["text"] = "\n\n".join(sheets_text)
                result["metadata"]["sheets"] = wb.sheetnames
                wb.close()
            except ImportError:
                result = self._extract_with_tika(file_path)
            except Exception as e:
                result["error"] = str(e)
        else:
            result = self._extract_with_tika(file_path)
        return result

    def _extract_powerpoint(self, file_path: str) -> dict[str, Any]:
        """Extract text from PowerPoint files."""
        result: dict[str, Any] = {"text": "", "method": "python-pptx", "metadata": {}, "error": ""}
        ext = Path(file_path).suffix.lower()
        if ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                slides_text = []
                for i, slide in enumerate(prs.slides, 1):
                    slide_lines = [f"[Slide {i}]"]
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                if para.text.strip():
                                    slide_lines.append(para.text)
                        if shape.has_table:
                            for row in shape.table.rows:
                                cells = [cell.text.strip() for cell in row.cells]
                                if any(cells):
                                    slide_lines.append(" | ".join(cells))
                    slides_text.append("\n".join(slide_lines))
                result["text"] = "\n\n".join(slides_text)
                result["metadata"]["slides"] = len(prs.slides)
            except ImportError:
                result = self._extract_with_tika(file_path)
            except Exception as e:
                result["error"] = str(e)
        else:
            result = self._extract_with_tika(file_path)
        return result

    def _extract_email(self, file_path: str) -> dict[str, Any]:
        """Extract text from .eml email files."""
        result: dict[str, Any] = {"text": "", "method": "email", "metadata": {}, "error": ""}
        try:
            with open(file_path, "rb") as f:
                msg = email.message_from_binary_file(f)
            parts = []
            result["metadata"] = {
                "from": msg.get("From", ""),
                "to": msg.get("To", ""),
                "cc": msg.get("Cc", ""),
                "subject": msg.get("Subject", ""),
                "date": msg.get("Date", ""),
            }
            parts.append(f"From: {msg.get('From', '')}")
            parts.append(f"To: {msg.get('To', '')}")
            parts.append(f"Subject: {msg.get('Subject', '')}")
            parts.append(f"Date: {msg.get('Date', '')}")
            parts.append("")

            attachments = []
            if msg.is_multipart():
                for part in msg.walk():
                    content_type = part.get_content_type()
                    disposition = str(part.get("Content-Disposition", ""))
                    if "attachment" in disposition:
                        fname = part.get_filename()
                        if fname:
                            attachments.append(fname)
                    elif content_type == "text/plain":
                        payload = part.get_payload(decode=True)
                        if payload:
                            parts.append(payload.decode("utf-8", errors="replace"))
                    elif content_type == "text/html":
                        payload = part.get_payload(decode=True)
                        if payload:
                            # Basic HTML stripping
                            import re
                            text = payload.decode("utf-8", errors="replace")
                            text = re.sub(r"<[^>]+>", " ", text)
                            text = re.sub(r"\s+", " ", text).strip()
                            parts.append(text)
            else:
                payload = msg.get_payload(decode=True)
                if payload:
                    parts.append(payload.decode("utf-8", errors="replace"))

            if attachments:
                parts.append(f"\nAttachments: {', '.join(attachments)}")
                result["metadata"]["attachments"] = attachments

            result["text"] = "\n".join(parts)
        except Exception as e:
            result["error"] = str(e)
        return result

    def _extract_msg(self, file_path: str) -> dict[str, Any]:
        """Extract text from Outlook .msg files."""
        result: dict[str, Any] = {"text": "", "method": "extract-msg", "metadata": {}, "error": ""}
        try:
            import extract_msg
            msg = extract_msg.Message(file_path)
            parts = [
                f"From: {msg.sender or ''}",
                f"To: {msg.to or ''}",
                f"Subject: {msg.subject or ''}",
                f"Date: {msg.date or ''}",
                "",
                msg.body or "",
            ]
            result["text"] = "\n".join(parts)
            result["metadata"] = {
                "from": msg.sender or "",
                "to": msg.to or "",
                "subject": msg.subject or "",
                "date": str(msg.date) if msg.date else "",
                "attachments": [a.longFilename for a in msg.attachments] if msg.attachments else [],
            }
            msg.close()
        except ImportError:
            result = self._extract_with_tika(file_path)
        except Exception as e:
            result["error"] = str(e)
        return result

    def _extract_image(self, file_path: str) -> dict[str, Any]:
        """Extract text from images via OCR, plus EXIF metadata."""
        result: dict[str, Any] = {"text": "", "method": "ocr", "metadata": {}, "error": ""}
        # EXIF metadata
        try:
            from PIL import Image
            from PIL.ExifTags import TAGS
            img = Image.open(file_path)
            result["metadata"]["size"] = f"{img.width}x{img.height}"
            result["metadata"]["format"] = img.format or ""
            exif = img._getexif()
            if exif:
                for tag_id, value in exif.items():
                    tag_name = TAGS.get(tag_id, str(tag_id))
                    if isinstance(value, (str, int, float)):
                        result["metadata"][tag_name] = value
            img.close()
        except Exception:
            pass

        # OCR
        try:
            import pytesseract
            from PIL import Image
            img = Image.open(file_path)
            text = pytesseract.image_to_string(img)
            result["text"] = text.strip()
            img.close()
        except ImportError:
            result["method"] = "exif-only"
            result["text"] = f"[Image: {result['metadata'].get('size', 'unknown')}]"
        except Exception as e:
            result["method"] = "exif-only"
            result["error"] = f"OCR failed: {e}"
        return result

    def _extract_media_metadata(self, file_path: str) -> dict[str, Any]:
        """Extract metadata from audio/video files."""
        result: dict[str, Any] = {"text": "", "method": "mutagen", "metadata": {}, "error": ""}
        try:
            from mutagen import File as MutagenFile
            audio = MutagenFile(file_path)
            if audio is not None:
                info_parts = []
                if audio.info:
                    if hasattr(audio.info, "length"):
                        mins = int(audio.info.length // 60)
                        secs = int(audio.info.length % 60)
                        info_parts.append(f"Duration: {mins}:{secs:02d}")
                        result["metadata"]["duration"] = audio.info.length
                    if hasattr(audio.info, "bitrate"):
                        result["metadata"]["bitrate"] = audio.info.bitrate
                if audio.tags:
                    for key in audio.tags:
                        val = audio.tags[key]
                        if isinstance(val, list):
                            val = ", ".join(str(v) for v in val)
                        info_parts.append(f"{key}: {val}")
                        result["metadata"][str(key)] = str(val)
                result["text"] = "\n".join(info_parts)
        except ImportError:
            result["error"] = "mutagen not installed"
        except Exception as e:
            result["error"] = str(e)
        return result

    def _extract_archive_listing(self, file_path: str) -> dict[str, Any]:
        """List contents of archive files."""
        result: dict[str, Any] = {"text": "", "method": "archive", "metadata": {}, "error": ""}
        ext = Path(file_path).suffix.lower()
        try:
            if ext == ".zip":
                with zipfile.ZipFile(file_path, "r") as zf:
                    names = zf.namelist()
                    result["text"] = f"Archive contents ({len(names)} files):\n" + "\n".join(names)
                    result["metadata"]["file_count"] = len(names)
                    result["metadata"]["files"] = names[:100]  # Cap for sanity
            else:
                result["text"] = f"[Archive: {ext}]"
                result["metadata"]["type"] = ext
        except Exception as e:
            result["error"] = str(e)
        return result

    def _extract_with_tika(self, file_path: str) -> dict[str, Any]:
        """Use Apache Tika for extraction (catch-all)."""
        result: dict[str, Any] = {"text": "", "method": "tika", "metadata": {}, "error": ""}
        try:
            from tika import parser as tika_parser
            if self.tika_server_url:
                parsed = tika_parser.from_file(file_path, serverEndpoint=self.tika_server_url)
            else:
                parsed = tika_parser.from_file(file_path)
            result["text"] = (parsed.get("content") or "").strip()
            result["metadata"] = parsed.get("metadata") or {}
        except ImportError:
            result["error"] = "Apache Tika (tika-python) not installed"
        except Exception as e:
            result["error"] = f"Tika extraction failed: {e}"
        return result
